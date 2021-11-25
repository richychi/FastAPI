from PIL import Image, ImageEnhance, ImageOps, ImageDraw, ImageFilter
import glob

from pptx import Presentation
from PIL.Image import image
import os

prs = Presentation()


class Slide001:
    def __init__(self, data):
        self.layout = prs.slide_layout[data[3]]
        self.slide = prs.slides.add_slide(self.layout)
        self.title = self.slide.shapes.title
        self.title.text = data[0]
        self.subtitle = self.slide.placeholders[1]
        self.subtitle.text = data[1]
        self.image = self.slide.placeholders[1].insert_picture(data[2])
        im = Image.open(self.image)


slides = [
    [
        "Title of the Presentation",
        "Subtitle of presentation",
        "image/OIP.png",
        0
    ],
    [
        "Title of paragraph",
        "Content (bullet)",
        "image/OIP2.png",
        1
    ],
    [
        "Another slide",
        "Content is this one that I'm writing just now.",
        "image/OIP.png",
        2
    ]
]

for es in slides:
    Slide001(es)

prs.save("MyPresentation.pptx")
os.startfile("MyPresentation.pptx")
