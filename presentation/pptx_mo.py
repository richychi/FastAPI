import pptx.util
from pptx import Presentation

prs = Presentation('fpptemplates/Template001.pptx')
# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]
#
# title.text = "Hello, World!"
# subtitle.text = "python-pptx was here!"
for slide in prs.slides:
    pic = slide.shapes.add_picture('image/IMT-Logo.png', pptx.util.Pt(550), pptx.util.Pt(50),
                                   width=pptx.util.Pt(100), height=pptx.util.Pt(100))

prs.save('image/test.pptx')
