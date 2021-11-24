# from PIL import pil

from pptx import Presentation

import os

prs = Presentation()


class Slide001:
    def __init__(self, data):
        self.layout = prs.slide_layout[data[3]]
        self.slide = prs.slides.add_slide(self.layout)
        self.title = self.slide.shapes.title
        self.title.text = data[0]
        self.subtitle = self.slide.placeholders[1]
        self.subtitle.text = data[1]
        self.image = self.slide.placeholders[1].insert_picture(data[2])


slides = [
    [
        "Title of the Presentation",
        "Subtitle of presentation",
        "005.png",
        0
    ],
    [
        "Title of paragraph",
        "Content (bullet)",
        "005.png",
        1
    ],
    [
        "Another slide",
        "Content is this one that I'm writing just now.",
        "005.png",
        2
    ]
]

for es in slides:
    Slide001(es)

prs.save("MyPresentation.pptx")
os.startfile("MyPresentation.pptx")
